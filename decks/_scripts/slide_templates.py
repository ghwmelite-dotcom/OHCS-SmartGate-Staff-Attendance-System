"""Slide factory functions — Kente Executive layouts with aspect-preserving image placement."""
from pathlib import Path
from PIL import Image
from pptx.util import Emu, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import theme

ASSETS = Path(__file__).parent.parent / "_assets"


# ─── Primitives ──────────────────────────────────────────────────────────────

def _set_slide_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, font, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def _add_picture_with_shadow(slide, image_path, left, top, width=None, height=None):
    """Add a picture and apply a soft drop-shadow via OOXML."""
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    spPr = pic._element.spPr
    # Remove existing effectLst if present
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        spPr.remove(existing)
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    effectLst = etree.SubElement(spPr, qn("a:effectLst"))
    outerShdw = etree.SubElement(
        effectLst, qn("a:outerShdw"),
        blurRad="120000",     # 0.13" blur
        dist="50000",         # 0.05" offset
        dir="5400000",        # 90° down
        algn="tl",
        rotWithShape="0",
    )
    srgbClr = etree.SubElement(outerShdw, qn("a:srgbClr"), val="000000")
    etree.SubElement(srgbClr, qn("a:alpha"), val="25000")  # 25% opacity
    return pic


def _fit_image_dims(image_path: str, max_w_emu: int, max_h_emu: int) -> tuple[int, int]:
    """Compute display dimensions that fit inside max_w x max_h while preserving native aspect ratio."""
    img = Image.open(image_path)
    src_w, src_h = img.size
    src_ar = src_w / src_h
    box_ar = max_w_emu / max_h_emu
    if src_ar >= box_ar:
        return max_w_emu, int(max_w_emu / src_ar)
    else:
        return int(max_h_emu * src_ar), max_h_emu


def _image_aspect_ratio(image_path: str) -> float:
    img = Image.open(image_path)
    return img.width / img.height


def _add_footer(slide, deck_id: str, deck_title: str, page: int, total: int):
    text = f"{deck_id}  ·  {deck_title}  ·  {page:02d} / {total:02d}"
    _add_text(
        slide,
        left=Emu(7500000), top=Emu(6500000),
        width=Emu(4500000), height=Emu(300000),
        text=text, font=theme.FONT_BODY, size=Pt(9),
        color=theme.NEUTRAL_LINE, align=PP_ALIGN.RIGHT,
    )


def _add_tier_strip(slide, tier: str):
    """Tier accent strip — 24px-tall top band (was 12px; doubled for visual presence)."""
    img_path = str(ASSETS / theme.TIER_STRIPS[tier])
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), width=theme.SLIDE_WIDTH_EMU, height=Emu(200000))


def _add_page_marker(slide, color=None):
    """Small gold tick at top-left of body slides — quiet, recurring brand mark."""
    if color is None:
        color = theme.GOLD_SIGNATURE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(150000), Emu(200000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ─── Cover slide ─────────────────────────────────────────────────────────────

def add_cover_slide(prs, deck_id, deck_title, deck_subtitle, tier, audience, version, date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, tier)

    # Eyebrow
    _add_text(
        slide, Emu(800000), Emu(800000), Emu(8000000), Emu(300000),
        f"OHCS  ·  {deck_id}", theme.FONT_BODY, Pt(11),
        theme.GOLD_DEEP, bold=True,
    )

    # Hero title — large, left-aligned for editorial weight
    _add_text(
        slide, Emu(800000), Emu(2300000), Emu(10500000), Emu(1800000),
        deck_title, theme.FONT_DISPLAY, Pt(72),
        theme.INK_DEEP, bold=True, align=PP_ALIGN.LEFT,
    )

    # Subtitle — set apart by whitespace, not a hairline
    _add_text(
        slide, Emu(800000), Emu(4400000), Emu(10500000), Emu(700000),
        deck_subtitle, theme.FONT_BODY, Pt(22), theme.INK_WARM, align=PP_ALIGN.LEFT,
    )

    # Bottom row
    _add_text(
        slide, Emu(800000), Emu(6200000), Emu(4000000), Emu(300000),
        f"For: {audience}", theme.FONT_BODY, Pt(11), theme.GOLD_DEEP, bold=True,
    )
    _add_text(
        slide, Emu(8000000), Emu(6200000), Emu(4000000), Emu(300000),
        f"v{version}  ·  {date}", theme.FONT_BODY, Pt(11), theme.NEUTRAL_LINE, align=PP_ALIGN.RIGHT,
    )
    return slide


# ─── Divider slide (opening / closing line) ──────────────────────────────────

def add_divider_slide(prs, line: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.INK_DEEP)
    slide.shapes.add_picture(
        str(ASSETS / "kente-texture-overlay.png"),
        Emu(0), Emu(0), width=theme.SLIDE_WIDTH_EMU, height=theme.SLIDE_HEIGHT_EMU,
    )
    _add_text(
        slide, Emu(1400000), Emu(2900000), Emu(9300000), Emu(2000000),
        line, theme.FONT_DISPLAY, Pt(40), theme.CREAM_PAGE,
        align=PP_ALIGN.CENTER,
    )
    return slide


# ─── TOC slide ───────────────────────────────────────────────────────────────

def add_toc_slide(prs, bullets: list[str], deck_id: str, deck_title: str, page: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "thematic")
    _add_text(
        slide, Emu(800000), Emu(800000), Emu(10000000), Emu(800000),
        "WHAT YOU'LL SEE", theme.FONT_BODY, Pt(11), theme.GOLD_DEEP, bold=True,
    )
    _add_text(
        slide, Emu(800000), Emu(1300000), Emu(10000000), Emu(800000),
        "Three movements.", theme.FONT_DISPLAY, Pt(48), theme.INK_DEEP, bold=True,
    )
    # Numbered list — generous spacing
    for i, bullet in enumerate(bullets):
        # Big numeral
        _add_text(
            slide, Emu(800000), Emu(2700000 + i * 1100000), Emu(700000), Emu(900000),
            f"{i+1:02d}", theme.FONT_DISPLAY, Pt(48), theme.GOLD_SIGNATURE, bold=True,
        )
        _add_text(
            slide, Emu(1700000), Emu(2900000 + i * 1100000), Emu(9500000), Emu(900000),
            bullet, theme.FONT_BODY, Pt(22), theme.INK_WARM,
        )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


# ─── Statement slide (text-only with subtle motif, no hairline) ──────────────

def add_statement_slide(prs, headline: str, sub: str, deck_id: str, deck_title: str, page: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "thematic")
    # Big gold quote mark as visual anchor (replaces the underline)
    _add_text(
        slide, Emu(800000), Emu(1100000), Emu(800000), Emu(800000),
        "//", theme.FONT_DISPLAY, Pt(72), theme.GOLD_SIGNATURE, bold=True,
    )
    # Headline — large, left-aligned
    _add_text(
        slide, Emu(800000), Emu(2200000), Emu(10500000), Emu(2200000),
        headline, theme.FONT_DISPLAY, Pt(48), theme.INK_DEEP, bold=True,
    )
    # Sub — gentle, distinct register
    _add_text(
        slide, Emu(800000), Emu(4700000), Emu(10500000), Emu(1500000),
        sub, theme.FONT_BODY, Pt(22), theme.INK_WARM,
    )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


# ─── Evidence slides — three layouts dispatched by aspect ratio ──────────────

def _evidence_landscape_wide(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total):
    """For AR >= 1.8 — image fills slide width, bullets sit below in 3 columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "smartgate")

    _add_text(
        slide, Emu(800000), Emu(500000), Emu(10500000), Emu(700000),
        title, theme.FONT_DISPLAY, Pt(28), theme.INK_DEEP, bold=True,
    )

    # Image — full width minus margins
    box_w = Emu(10592000)
    box_h = Emu(3600000)
    w, h = _fit_image_dims(image_path, box_w, box_h)
    img_left = Emu(800000) + (box_w - w) // 2
    img_top = Emu(1500000) + (box_h - h) // 2
    _add_picture_with_shadow(slide, image_path, img_left, img_top, width=w, height=h)

    # Bullets — three columns below the image
    col_w = Emu(3400000)
    col_gap = Emu(196000)
    base_left = Emu(800000)
    bullets_top = Emu(5400000)
    for i, bullet in enumerate(bullets[:3]):
        # Numeral
        _add_text(
            slide,
            base_left + i * (col_w + col_gap),
            bullets_top,
            Emu(400000), Emu(400000),
            f"{i+1:02d}", theme.FONT_DISPLAY, Pt(20), theme.GOLD_SIGNATURE, bold=True,
        )
        _add_text(
            slide,
            base_left + i * (col_w + col_gap),
            bullets_top + Emu(450000),
            col_w, Emu(900000),
            bullet, theme.FONT_BODY, Pt(13), theme.INK_WARM,
        )

    _add_text(
        slide, Emu(800000), Emu(6480000), Emu(10500000), Emu(250000),
        caption, theme.FONT_BODY, Pt(10), theme.GOLD_DEEP,
    )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


def _evidence_landscape_standard(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total):
    """For 1.15 <= AR < 1.8 — image dominant left (anchored left, not centered), bullets right column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "smartgate")

    _add_text(
        slide, Emu(800000), Emu(500000), Emu(10500000), Emu(700000),
        title, theme.FONT_DISPLAY, Pt(28), theme.INK_DEEP, bold=True,
    )

    # Image — anchored to LEFT margin (no horizontal centering inside box)
    box_w = Emu(7400000)
    box_h = Emu(4900000)
    w, h = _fit_image_dims(image_path, box_w, box_h)
    img_left = Emu(800000)
    img_top = Emu(1500000) + (box_h - h) // 2  # vertical center is fine
    _add_picture_with_shadow(slide, image_path, img_left, img_top, width=w, height=h)

    # Bullets — right column starts just past where image ends (or at fixed 8.4M if shorter)
    bx = max(img_left + w + Emu(500000), Emu(8400000))
    bw = theme.SLIDE_WIDTH_EMU - bx - Emu(500000)
    for i, bullet in enumerate(bullets[:3]):
        _add_text(
            slide, bx, Emu(1700000) + i * Emu(1400000),
            Emu(400000), Emu(400000),
            f"{i+1:02d}", theme.FONT_DISPLAY, Pt(22), theme.GOLD_SIGNATURE, bold=True,
        )
        _add_text(
            slide, bx, Emu(2100000) + i * Emu(1400000),
            bw, Emu(1200000),
            bullet, theme.FONT_BODY, Pt(15), theme.INK_WARM,
        )

    _add_text(
        slide, Emu(800000), Emu(6480000), Emu(10500000), Emu(250000),
        caption, theme.FONT_BODY, Pt(10), theme.GOLD_DEEP,
    )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


def _evidence_square(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total):
    """For 0.85 <= AR < 1.15 — near-square image, image left-anchored at large size, bullets in narrow right column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "smartgate")

    _add_text(
        slide, Emu(800000), Emu(500000), Emu(10500000), Emu(700000),
        title, theme.FONT_DISPLAY, Pt(28), theme.INK_DEEP, bold=True,
    )

    # Image — square-ish; size to fit max height (5M EMU), then place to left
    box_h = Emu(5000000)
    w, h = _fit_image_dims(image_path, Emu(7500000), box_h)
    img_left = Emu(800000)
    img_top = Emu(1500000)
    _add_picture_with_shadow(slide, image_path, img_left, img_top, width=w, height=h)

    # Bullets — narrow right column, vertically centered against image
    bx = img_left + w + Emu(500000)
    bw = theme.SLIDE_WIDTH_EMU - bx - Emu(500000)
    if bw < Emu(2200000):
        bw = Emu(2200000)
        bx = theme.SLIDE_WIDTH_EMU - bw - Emu(500000)
    bullets_top = img_top + (h - Emu(3 * 1100000)) // 2  # center against image height
    for i, bullet in enumerate(bullets[:3]):
        _add_text(
            slide, bx, bullets_top + i * Emu(1100000),
            Emu(400000), Emu(400000),
            f"{i+1:02d}", theme.FONT_DISPLAY, Pt(20), theme.GOLD_SIGNATURE, bold=True,
        )
        _add_text(
            slide, bx, bullets_top + Emu(400000) + i * Emu(1100000),
            bw, Emu(900000),
            bullet, theme.FONT_BODY, Pt(13), theme.INK_WARM,
        )

    _add_text(
        slide, Emu(800000), Emu(6700000), Emu(10500000), Emu(250000),
        caption, theme.FONT_BODY, Pt(10), theme.GOLD_DEEP,
    )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


def _evidence_portrait(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total):
    """For AR < 0.85 — true portrait (mobile-like); image right, bullets left, image larger and more dramatic."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "smartgate")

    _add_text(
        slide, Emu(800000), Emu(500000), Emu(10500000), Emu(700000),
        title, theme.FONT_DISPLAY, Pt(28), theme.INK_DEEP, bold=True,
    )

    # Image — true portrait, generous size on the right
    box_w = Emu(4400000)
    box_h = Emu(5200000)
    w, h = _fit_image_dims(image_path, box_w, box_h)
    img_left = Emu(7400000) + (box_w - w) // 2
    img_top = Emu(1400000) + (box_h - h) // 2
    _add_picture_with_shadow(slide, image_path, img_left, img_top, width=w, height=h)

    # Bullets — left column, larger type
    bx = Emu(800000)
    bw = Emu(6200000)
    bullets_top = img_top + Emu(300000)
    for i, bullet in enumerate(bullets[:3]):
        _add_text(
            slide, bx, bullets_top + i * Emu(1500000),
            Emu(500000), Emu(500000),
            f"{i+1:02d}", theme.FONT_DISPLAY, Pt(28), theme.GOLD_SIGNATURE, bold=True,
        )
        _add_text(
            slide, bx, bullets_top + Emu(500000) + i * Emu(1500000),
            bw, Emu(1300000),
            bullet, theme.FONT_BODY, Pt(17), theme.INK_WARM,
        )

    _add_text(
        slide, Emu(800000), Emu(6700000), Emu(10500000), Emu(250000),
        caption, theme.FONT_BODY, Pt(10), theme.GOLD_DEEP,
    )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


def add_evidence_slide(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total):
    """Dispatch to the right layout based on image aspect ratio.

    Bands:  AR >= 1.8  → ultra-wide (image fills width, bullets below in 3 cols)
            1.15 <= AR < 1.8 → standard landscape (image left, bullets right)
            0.85 <= AR < 1.15 → square (image left at full height, bullets vertically centered right)
            AR < 0.85 → portrait (image right at larger size, bullets left)
    """
    if image_path and Path(image_path).exists():
        ar = _image_aspect_ratio(image_path)
        if ar >= 1.8:
            return _evidence_landscape_wide(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total)
        elif ar >= 1.15:
            return _evidence_landscape_standard(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total)
        elif ar >= 0.85:
            return _evidence_square(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total)
        else:
            return _evidence_portrait(prs, title, image_path, bullets, caption, deck_id, deck_title, page, total)
    # Fallback — placeholder slide when no image yet
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "smartgate")
    _add_text(slide, Emu(800000), Emu(500000), Emu(10500000), Emu(700000),
              title, theme.FONT_DISPLAY, Pt(28), theme.INK_DEEP, bold=True)
    ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Emu(800000), Emu(1500000), Emu(10500000), Emu(4500000))
    ph.fill.solid()
    ph.fill.fore_color.rgb = theme.CREAM_SOFT
    ph.line.color.rgb = theme.NEUTRAL_LINE
    tf = ph.text_frame
    tf.text = f"[REPLACE: {Path(image_path).name if image_path else 'image'}]"
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = theme.FONT_BODY
            r.font.size = Pt(14)
            r.font.color.rgb = theme.GOLD_DEEP
    _add_text(slide, Emu(800000), Emu(6480000), Emu(10500000), Emu(250000),
              caption, theme.FONT_BODY, Pt(10), theme.GOLD_DEEP)
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


# ─── Wow slide (hero number) ─────────────────────────────────────────────────

def add_wow_slide(prs, hero_number: str, label: str, deck_id: str, deck_title: str, page: int, total: int):
    """Hero number slide. Gracefully handles <NXX> placeholder tokens by rendering a
    composed 'pending figure' treatment instead of literal angle brackets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.INK_DEEP)
    slide.shapes.add_picture(
        str(ASSETS / "kente-texture-overlay.png"),
        Emu(0), Emu(0), width=theme.SLIDE_WIDTH_EMU, height=theme.SLIDE_HEIGHT_EMU,
    )

    is_placeholder = (
        isinstance(hero_number, str)
        and hero_number.startswith("<")
        and hero_number.endswith(">")
    )

    if is_placeholder:
        # Render an intentional "figure pending" treatment so the slide doesn't
        # look like an exposed token. Designed to look complete on its own.
        _add_text(
            slide, Emu(800000), Emu(1200000), Emu(10500000), Emu(400000),
            "THE NUMBER", theme.FONT_BODY, Pt(12), theme.GOLD_SIGNATURE, bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, Emu(800000), Emu(2100000), Emu(10500000), Emu(3300000),
            "—", theme.FONT_DISPLAY, Pt(260),
            theme.GOLD_SIGNATURE, bold=True, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, Emu(1800000), Emu(5500000), Emu(8500000), Emu(700000),
            label, theme.FONT_BODY, Pt(22), theme.CREAM_PAGE, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, Emu(800000), Emu(6300000), Emu(10500000), Emu(300000),
            f"Figure pending  ·  fill {hero_number} in outline before delivery",
            theme.FONT_BODY, Pt(10), theme.NEUTRAL_LINE, align=PP_ALIGN.CENTER,
        )
    else:
        _add_text(
            slide, Emu(800000), Emu(1200000), Emu(10500000), Emu(400000),
            "THE NUMBER", theme.FONT_BODY, Pt(12), theme.GOLD_SIGNATURE, bold=True,
            align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, Emu(800000), Emu(2100000), Emu(10500000), Emu(3300000),
            hero_number, theme.FONT_DISPLAY, Pt(260),
            theme.GOLD_SIGNATURE, bold=True, align=PP_ALIGN.CENTER,
        )
        _add_text(
            slide, Emu(1800000), Emu(5500000), Emu(8500000), Emu(700000),
            label, theme.FONT_BODY, Pt(22), theme.CREAM_PAGE, align=PP_ALIGN.CENTER,
        )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide


# ─── Appendix slide ──────────────────────────────────────────────────────────

def add_appendix_slide(prs, links: list[str], related: list[str], deck_id: str, deck_title: str, page: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, theme.CREAM_PAGE)
    _add_tier_strip(slide, "thematic")

    _add_text(
        slide, Emu(800000), Emu(800000), Emu(10000000), Emu(400000),
        "APPENDIX", theme.FONT_BODY, Pt(11), theme.GOLD_DEEP, bold=True,
    )
    _add_text(
        slide, Emu(800000), Emu(1300000), Emu(10000000), Emu(800000),
        "Where to look next.", theme.FONT_DISPLAY, Pt(40), theme.INK_DEEP, bold=True,
    )

    # Two columns
    _add_text(
        slide, Emu(800000), Emu(2900000), Emu(5000000), Emu(400000),
        "LIVE", theme.FONT_BODY, Pt(11), theme.GOLD_DEEP, bold=True,
    )
    for i, link in enumerate(links):
        _add_text(slide, Emu(800000), Emu(3400000 + i * 480000), Emu(5500000), Emu(400000),
                  link, theme.FONT_BODY, Pt(15), theme.INK_WARM)

    _add_text(
        slide, Emu(6800000), Emu(2900000), Emu(5000000), Emu(400000),
        "COMPANION", theme.FONT_BODY, Pt(11), theme.GOLD_DEEP, bold=True,
    )
    for i, rel in enumerate(related):
        _add_text(slide, Emu(6800000), Emu(3400000 + i * 480000), Emu(5000000), Emu(400000),
                  rel, theme.FONT_BODY, Pt(15), theme.INK_WARM)

    _add_text(
        slide, Emu(800000), Emu(6300000), Emu(10000000), Emu(400000),
        "Built on Cloudflare's edge  ·  Designed in the Kente Executive language",
        theme.FONT_BODY, Pt(10), theme.GOLD_DEEP, align=PP_ALIGN.LEFT,
    )
    _add_footer(slide, deck_id, deck_title, page, total)
    return slide
