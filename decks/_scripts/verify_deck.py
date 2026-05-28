"""Reads back a .pptx and asserts: slide count matches, file size under cap, no broken pictures."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def verify(pptx_path: Path, expected_slides: int, max_mb: int = 25) -> None:
    prs = Presentation(str(pptx_path))
    actual = len(prs.slides)
    if actual != expected_slides:
        raise AssertionError(f"{pptx_path.name}: expected {expected_slides} slides, got {actual}")
    size_mb = pptx_path.stat().st_size / (1024 * 1024)
    if size_mb > max_mb:
        raise AssertionError(f"{pptx_path.name}: {size_mb:.1f} MB exceeds {max_mb} MB cap")
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.image is None:
                    raise AssertionError(f"{pptx_path.name}: slide {i} has orphan picture")
    print(f"OK  {pptx_path.name}  slides={actual}  size={size_mb:.2f}MB")


if __name__ == "__main__":
    pptx_path = Path(sys.argv[1])
    expected = int(sys.argv[2])
    verify(pptx_path, expected)
