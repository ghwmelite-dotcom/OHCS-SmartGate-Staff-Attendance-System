"""Three style mockups of the same slide (Check-In Step 1) for the user to compare.

Renders each mockup as a 1920x1080 PNG via Pillow, then bundles them into a
3-slide comparison .pptx. Each mockup is embedded full-bleed so the user is
seeing exactly what the final approach would produce.
"""
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Emu

# ─── Paths and constants ─────────────────────────────────────────────────────

ROOT = Path(__file__).parent.parent
ASSETS = ROOT / "_assets"
SCREENS = ASSETS / "screenshots" / "OHCS VMS"
OUT_DIR = ROOT / "output"
MOCKUP_DIR = ASSETS / "mockups"
MOCKUP_DIR.mkdir(parents=True, exist_ok=True)

SCREENSHOT = SCREENS / "OHCS-VMS-—-Check_In_p1.png"

W, H = 1920, 1080
WIN = "C:/Windows/Fonts"

# Font paths
F_SERIF_BOLD     = f"{WIN}/georgiab.ttf"
F_SERIF_REGULAR  = f"{WIN}/georgia.ttf"
F_SERIF_ITALIC   = f"{WIN}/georgiai.ttf"
F_SANS_REGULAR   = f"{WIN}/segoeui.ttf"
F_SANS_BOLD      = f"{WIN}/segoeuib.ttf"
F_SANS_BLACK     = f"{WIN}/seguibl.ttf"

# Palette
INK_DEEP        = (0x0E, 0x14, 0x11)
INK_WARM        = (0x1A, 0x17, 0x14)
CREAM_PAGE      = (0xF6, 0xF1, 0xE7)
CREAM_SOFT      = (0xFB, 0xF7, 0xEF)
GOLD_SIGNATURE  = (0xC9, 0xA1, 0x4A)
GOLD_DEEP       = (0x8B, 0x6B, 0x22)
GREEN_SMARTGATE = (0x1A, 0x4D, 0x2E)
NEUTRAL_LINE    = (0xD8, 0xCF, 0xBE)
WHITE           = (0xFF, 0xFF, 0xFF)
NAVY            = (0x1E, 0x27, 0x47)


def ft(path: str, size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(path, size)


def load_screenshot() -> Image.Image:
    return Image.open(SCREENSHOT).convert("RGB")


def add_drop_shadow(img: Image.Image, offset=(0, 24), blur=40, opacity=120) -> Image.Image:
    """Returns an RGBA image with a soft drop-shadow composited beneath the original."""
    w, h = img.size
    pad = blur + max(abs(offset[0]), abs(offset[1])) + 10
    canvas = Image.new("RGBA", (w + pad * 2, h + pad * 2), (0, 0, 0, 0))
    shadow = Image.new("RGBA", (w, h), (0, 0, 0, opacity))
    canvas.paste(shadow, (pad + offset[0], pad + offset[1]), shadow)
    canvas = canvas.filter(ImageFilter.GaussianBlur(blur))
    canvas.paste(img.convert("RGBA"), (pad, pad), img.convert("RGBA") if img.mode == "RGBA" else None)
    return canvas


def draw_browser_chrome(canvas: Image.Image, x: int, y: int, w: int, h: int, url: str, chrome_h: int = 56) -> tuple[int, int, int, int]:
    """Draws a macOS-style browser chrome on top of the canvas at (x,y) of size w×h.
    Returns the (left, top, right, bottom) box where the screenshot should sit beneath the chrome.
    """
    d = ImageDraw.Draw(canvas)
    # Outer rounded rect with white chrome
    chrome_box = (x, y, x + w, y + chrome_h)
    d.rounded_rectangle((x, y, x + w, y + h), radius=12, fill=(245, 244, 240))
    # Chrome top band
    d.rectangle(chrome_box, fill=(232, 230, 224))
    # Traffic lights
    cy = y + chrome_h // 2
    for i, color in enumerate([(252, 96, 89), (252, 188, 47), (40, 200, 64)]):
        cx = x + 22 + i * 22
        d.ellipse((cx - 7, cy - 7, cx + 7, cy + 7), fill=color)
    # URL pill
    url_x1 = x + 110
    url_x2 = x + w - 60
    url_y1 = y + 12
    url_y2 = y + chrome_h - 12
    d.rounded_rectangle((url_x1, url_y1, url_x2, url_y2), radius=8, fill=(220, 218, 210))
    url_font = ft(F_SANS_REGULAR, 16)
    d.text((url_x1 + 18, url_y1 + 6), url, font=url_font, fill=(80, 80, 80))
    # Return the inner content box for the screenshot
    return (x, y + chrome_h, x + w, y + h)


def fit_image_to_box(img: Image.Image, box_w: int, box_h: int) -> Image.Image:
    """Resize image preserving aspect ratio to fit box."""
    src_w, src_h = img.size
    src_ar = src_w / src_h
    box_ar = box_w / box_h
    if src_ar >= box_ar:
        new_w = box_w
        new_h = int(box_w / src_ar)
    else:
        new_h = box_h
        new_w = int(box_h * src_ar)
    return img.resize((new_w, new_h), Image.LANCZOS)


# ─── MOCKUP A — Dark editorial (Apple keynote / Stripe) ──────────────────────

def render_dark() -> Path:
    canvas = Image.new("RGB", (W, H), INK_DEEP)
    d = ImageDraw.Draw(canvas)

    # Subtle gold thread top-left as brand mark
    d.rectangle((0, 0, 6, 80), fill=GOLD_SIGNATURE)

    # Eyebrow
    eyebrow_font = ft(F_SANS_BOLD, 17)
    d.text((80, 60), "OHCS / SMARTGATE — CHECK-IN FLOW", font=eyebrow_font, fill=GOLD_SIGNATURE)

    # Screenshot in a macOS browser frame, centered
    src = load_screenshot()
    frame_w = 1400
    src_fit = fit_image_to_box(src, frame_w - 40, 720)
    frame_h = src_fit.height + 56 + 20  # chrome + screenshot + padding
    frame_x = (W - frame_w) // 2
    frame_y = 200

    # Build the browser frame on a transparent canvas, then drop-shadow + paste
    frame_canvas = Image.new("RGBA", (frame_w, frame_h), (0, 0, 0, 0))
    inner = draw_browser_chrome(frame_canvas, 0, 0, frame_w, frame_h, "smartgate.ohcsghana.org/check-in")
    ix, iy, ix2, iy2 = inner
    box_w = ix2 - ix
    sx = ix + (box_w - src_fit.width) // 2
    frame_canvas.paste(src_fit, (sx, iy + 10))

    # Drop-shadow the entire frame
    shadow_canvas = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    shadow_block = Image.new("RGBA", (frame_w, frame_h), (0, 0, 0, 160))
    shadow_canvas.paste(shadow_block, (frame_x + 0, frame_y + 30))
    shadow_canvas = shadow_canvas.filter(ImageFilter.GaussianBlur(35))
    canvas = Image.alpha_composite(canvas.convert("RGBA"), shadow_canvas).convert("RGB")
    canvas.paste(frame_canvas, (frame_x, frame_y), frame_canvas)

    d = ImageDraw.Draw(canvas)
    # Headline
    title_font = ft(F_SERIF_BOLD, 64)
    d.text((80, H - 240), "Step 01  Find or register", font=title_font, fill=CREAM_PAGE)

    # Sub
    sub_font = ft(F_SANS_REGULAR, 22)
    d.text((82, H - 150), "Existing visitor? Search. New visitor? Three fields. Search is type-ahead — no waiting.",
           font=sub_font, fill=NEUTRAL_LINE)

    # Page number
    pg_font = ft(F_SANS_REGULAR, 13)
    d.text((W - 130, H - 50), "06 / 28", font=pg_font, fill=NEUTRAL_LINE)

    out = MOCKUP_DIR / "mockup_a_dark_editorial.png"
    canvas.save(out)
    print(f"  rendered {out.name}")
    return out


# ─── MOCKUP B — Magazine spread (Monocle / WIRED) ────────────────────────────

def render_magazine() -> Path:
    canvas = Image.new("RGB", (W, H), CREAM_SOFT)
    d = ImageDraw.Draw(canvas)

    # Thin gold rule down the left edge
    d.rectangle((0, 0, 3, H), fill=GOLD_SIGNATURE)

    # Eyebrow
    eyebrow_font = ft(F_SANS_BOLD, 16)
    d.text((90, 80), "S T E P   O N E", font=eyebrow_font, fill=GOLD_DEEP)

    # HUGE serif title
    title_font = ft(F_SERIF_BOLD, 96)
    d.text((84, 130), "Find or register", font=title_font, fill=INK_DEEP)
    d.text((84, 245), "the visitor.", font=title_font, fill=INK_DEEP)

    # Italic subtitle
    sub_font = ft(F_SERIF_ITALIC, 28)
    d.text((90, 380), "Existing? Search.  New? Three fields.", font=sub_font, fill=INK_WARM)

    # Screenshot on the right at ~52% width
    src = load_screenshot()
    img_w = 900
    src_fit = fit_image_to_box(src, img_w, 680)
    img_x = W - src_fit.width - 90
    img_y = (H - src_fit.height) // 2 + 40

    # Subtle shadow
    shadow_canvas = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    shadow_block = Image.new("RGBA", (src_fit.width, src_fit.height), (0, 0, 0, 80))
    shadow_canvas.paste(shadow_block, (img_x + 8, img_y + 16))
    shadow_canvas = shadow_canvas.filter(ImageFilter.GaussianBlur(18))
    canvas = Image.alpha_composite(canvas.convert("RGBA"), shadow_canvas).convert("RGB")

    canvas.paste(src_fit, (img_x, img_y))

    # Thin gold hairline above image for editorial framing
    d = ImageDraw.Draw(canvas)
    d.rectangle((img_x, img_y - 14, img_x + 80, img_y - 11), fill=GOLD_SIGNATURE)

    # Pull-quote box (left mid)
    pq_y = 520
    d.rectangle((84, pq_y, 84 + 60, pq_y + 4), fill=GOLD_SIGNATURE)
    quote_font = ft(F_SERIF_ITALIC, 26)
    quote_lines = [
        "“Reception didn't have to learn it.",
        "They already knew it.”",
    ]
    for i, line in enumerate(quote_lines):
        d.text((84, pq_y + 30 + i * 40), line, font=quote_font, fill=INK_WARM)

    attribution = ft(F_SANS_BOLD, 13)
    d.text((84, pq_y + 130), "RECEPTION LEAD, OHCS HQ", font=attribution, fill=GOLD_DEEP)

    # Stat callout below pull-quote
    stat_font = ft(F_SERIF_BOLD, 78)
    d.text((84, pq_y + 200), "< 60s", font=stat_font, fill=GREEN_SMARTGATE)
    stat_label_font = ft(F_SANS_REGULAR, 16)
    d.text((84, pq_y + 295), "Average time from visitor arrival to printed badge.",
           font=stat_label_font, fill=INK_WARM)

    # Caption + page bottom
    cap_font = ft(F_SANS_REGULAR, 12)
    d.text((84, H - 50), "C H E C K - I N   F L O W   ·   O P E N I N G   S T A T E", font=cap_font, fill=GOLD_DEEP)
    d.text((W - 140, H - 50), "06 OF 28", font=cap_font, fill=GOLD_DEEP)

    out = MOCKUP_DIR / "mockup_b_magazine_spread.png"
    canvas.save(out)
    print(f"  rendered {out.name}")
    return out


# ─── MOCKUP C — Premium consulting (McKinsey / BCG) ──────────────────────────

def render_consulting() -> Path:
    canvas = Image.new("RGB", (W, H), WHITE)
    d = ImageDraw.Draw(canvas)

    # Top header bar
    d.rectangle((0, 0, W, 56), fill=NAVY)
    hdr_font = ft(F_SANS_BOLD, 13)
    d.text((80, 21), "OHCS  ·  SMARTGATE  ·  VISITOR MANAGEMENT  ·  EXECUTIVE BRIEFING", font=hdr_font, fill=WHITE)
    d.text((W - 180, 21), "06  /  28", font=hdr_font, fill=NEUTRAL_LINE)

    # Title block
    section_font = ft(F_SANS_BOLD, 14)
    d.text((80, 100), "1.  RECEPTION WORKFLOW", font=section_font, fill=GOLD_DEEP)

    title_font = ft(F_SANS_BOLD, 38)
    d.text((80, 132), "Step 1 of 4 — Find or register the visitor", font=title_font, fill=INK_DEEP)

    # Underline below the title (thin)
    d.rectangle((80, 195, 760, 197), fill=NAVY)

    # Source kicker
    src_font = ft(F_SANS_REGULAR, 12)
    d.text((80, 210), "Source: smartgate.ohcsghana.org/check-in  ·  Captured May 2026", font=src_font, fill=(120, 120, 120))

    # Main content split: image left (55%), right column key takeaways
    src = load_screenshot()
    img_w, img_h = 1050, 700
    src_fit = fit_image_to_box(src, img_w, img_h)
    img_x = 80
    img_y = 260
    # Thin border
    d.rectangle((img_x - 1, img_y - 1, img_x + src_fit.width + 1, img_y + src_fit.height + 1), outline=(200, 200, 200), width=1)
    canvas.paste(src_fit, (img_x, img_y))

    # Right column — Key Takeaways
    col_x = 1200
    d.text((col_x, 260), "KEY TAKEAWAYS", font=section_font, fill=GOLD_DEEP)

    takeaways = [
        ("Search-first design.", "Existing visitors are surfaced in two characters; new visitors require three fields."),
        ("Single workflow.", "Reception completes check-in without leaving this screen — host routing and purpose tagging inline."),
        ("Sub-60-second time-to-badge.", "Internal benchmark; consistent across reception staff with no training overhead."),
    ]

    bullet_num_font = ft(F_SANS_BLACK, 22)
    bullet_head_font = ft(F_SANS_BOLD, 16)
    bullet_body_font = ft(F_SANS_REGULAR, 13)

    for i, (head, body) in enumerate(takeaways):
        y0 = 310 + i * 160
        d.text((col_x, y0), f"{i+1:02d}", font=bullet_num_font, fill=NAVY)
        d.text((col_x + 44, y0 + 6), head, font=bullet_head_font, fill=INK_DEEP)
        # Wrap body manually (simple)
        body_lines = wrap_text(body, bullet_body_font, 580)
        for j, line in enumerate(body_lines):
            d.text((col_x + 44, y0 + 36 + j * 22), line, font=bullet_body_font, fill=(80, 80, 80))

    # Bottom footer
    d.rectangle((0, H - 36, W, H - 35), fill=(220, 220, 220))
    foot_font = ft(F_SANS_REGULAR, 11)
    d.text((80, H - 28), "OHCS SmartGate Executive Briefing  ·  v0  ·  28 May 2026", font=foot_font, fill=(140, 140, 140))
    d.text((W - 230, H - 28), "Confidential — for internal review only", font=foot_font, fill=(140, 140, 140))

    out = MOCKUP_DIR / "mockup_c_premium_consulting.png"
    canvas.save(out)
    print(f"  rendered {out.name}")
    return out


def wrap_text(text: str, font: ImageFont.FreeTypeFont, max_w: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    line = ""
    for w in words:
        candidate = (line + " " + w).strip()
        bbox = font.getbbox(candidate)
        if bbox[2] - bbox[0] <= max_w:
            line = candidate
        else:
            if line:
                lines.append(line)
            line = w
    if line:
        lines.append(line)
    return lines


# ─── Bundle into a comparison PPTX ───────────────────────────────────────────

def build_comparison_pptx(mockups: list[Path]) -> Path:
    prs = Presentation()
    # 16:9
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank = prs.slide_layouts[6]
    for png in mockups:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)
    out = OUT_DIR / "MOCKUP-COMPARISON.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    print("Rendering 3 mockup styles for the Check-In Step 1 slide...")
    a = render_dark()
    b = render_magazine()
    c = render_consulting()
    out = build_comparison_pptx([a, b, c])
    print(f"\nComparison deck: {out}")
    print(f"  Slide 1 — Dark editorial (Apple keynote / Stripe launch)")
    print(f"  Slide 2 — Magazine spread (Monocle / WIRED)")
    print(f"  Slide 3 — Premium consulting (McKinsey / BCG)")
    print(f"\nAlso saved individual PNGs to {MOCKUP_DIR}/")
