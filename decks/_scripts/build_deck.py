"""Orchestrates a deck build: outline.md -> .pptx via the magazine renderer.

Each slide is rendered as a high-res PNG (via magazine_renderer) and embedded
full-bleed in the .pptx. This gives complete design control — every slide is
essentially a designed image, the way premium product-launch decks are built.

Trade-off: text inside slides is NOT editable in PowerPoint. To change wording,
edit the outline .md and rebuild.
"""
from pathlib import Path
from tempfile import TemporaryDirectory
from pptx import Presentation

import theme
import magazine_renderer as mag
from outline_loader import parse_outline


def build(outline_path: Path, output_path: Path, deck_id: str, deck_title: str):
    slides = parse_outline(outline_path)
    if not slides:
        raise ValueError(f"No slides parsed from {outline_path}")

    prs = Presentation()
    prs.slide_width = theme.SLIDE_WIDTH_EMU
    prs.slide_height = theme.SLIDE_HEIGHT_EMU
    blank_layout = prs.slide_layouts[6]

    total = len(slides)
    section_eyebrows = _derive_eyebrows(slides, deck_title)

    with TemporaryDirectory() as td:
        td_path = Path(td)
        for i, slide_data in enumerate(slides, start=1):
            t = slide_data["type"]
            eyebrow = section_eyebrows.get(i, "")
            img = _render_one(
                t, slide_data, deck_id, deck_title, i, total, eyebrow,
                outline_path.parent, slides,
            )
            png_path = td_path / f"slide_{i:02d}.png"
            img.save(png_path, optimize=True)
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(png_path), 0, 0,
                                     width=prs.slide_width, height=prs.slide_height)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
    return output_path


def _derive_eyebrows(slides: list[dict], deck_title: str) -> dict[int, str]:
    """Walk the outline and assign a 'section eyebrow' to each slide based on the most
    recent preceding statement headline (or fallback to the deck title)."""
    out: dict[int, str] = {}
    current = deck_title.split("—")[0].strip()
    for i, sd in enumerate(slides, start=1):
        if sd["type"] == "statement":
            current = sd.get("headline", current).split(".")[0].strip()
            if len(current) > 38:
                current = current[:35].rstrip() + "…"
        out[i] = current
    return out


def _is_section_opener(slides: list[dict], idx: int) -> bool:
    """A statement slide is a section opener iff the next slide is evidence or wow."""
    if idx + 1 >= len(slides):
        return False
    return slides[idx + 1]["type"] in ("evidence", "wow")


def _statement_part_num(slides: list[dict], idx: int) -> int:
    """1-based count of section-opener statements up to and including this one."""
    n = 0
    for j in range(idx + 1):
        if slides[j]["type"] == "statement" and _is_section_opener(slides, j):
            n += 1
    return n


def _section_preview(slides: list[dict], idx: int) -> list[str]:
    """Titles of slides following idx until the next statement / divider / end."""
    titles: list[str] = []
    j = idx + 1
    while j < len(slides) and slides[j]["type"] not in ("statement", "divider"):
        sd = slides[j]
        if sd["type"] == "evidence":
            titles.append(sd.get("title", ""))
        elif sd["type"] == "wow":
            titles.append(f"The number — {sd.get('label', '')}")
        j += 1
    return titles[:5]


def _render_one(t: str, sd: dict, deck_id: str, deck_title: str, page: int, total: int,
                eyebrow: str, outline_dir: Path, all_slides: list[dict]):
    if t == "cover":
        return mag.render_cover(
            deck_id=deck_id,
            title=sd["title"],
            subtitle=sd.get("subtitle", ""),
            audience=sd.get("audience", ""),
            version=sd.get("version", "0"),
            date=sd.get("date", ""),
        )
    if t == "divider":
        return mag.render_divider(sd["line"])
    if t == "toc":
        return mag.render_toc(sd.get("bullets", []), page, total)
    if t == "statement":
        idx = page - 1
        if _is_section_opener(all_slides, idx):
            return mag.render_statement(
                headline=sd["headline"],
                sub=sd.get("sub", ""),
                page=page, total=total,
                part_num=_statement_part_num(all_slides, idx),
                preview_titles=_section_preview(all_slides, idx),
            )
        # Standalone statement (e.g., quote slide) — no part chrome
        return mag.render_statement(
            headline=sd["headline"],
            sub=sd.get("sub", ""),
            page=page, total=total,
            eyebrow=deck_title.split("—")[0].strip(),
        )
    if t == "evidence":
        image = sd.get("image")
        image_full = str((outline_dir / image).resolve()) if image else None
        return mag.render_evidence(
            title=sd["title"],
            image_path=image_full,
            bullets=sd.get("bullets", []),
            caption=sd.get("caption", ""),
            page=page, total=total,
            eyebrow=eyebrow,
        )
    if t == "wow":
        return mag.render_wow(sd["hero"], sd.get("label", ""), page, total)
    if t == "appendix":
        return mag.render_appendix(sd.get("links", []), sd.get("related", []), page, total)
    raise ValueError(f"Unknown slide type: {t}")
